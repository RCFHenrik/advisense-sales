"""AI-powered presentation analysis for campaign contact suggestions.

Extracts text from PPTX/PDF files, uses Anthropic Claude API to identify
business themes, and matches those themes against contact RelevanceTags.
"""

import json
import os
import logging
from typing import List

import httpx
from sqlalchemy.orm import Session

from app.models.models import Contact, ContactStatusEnum, SuppressionEntry

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from PyPDF2 import PdfReader
    HAS_PDF = True
except ImportError:
    HAS_PDF = False


class PresentationAnalysisService:
    def __init__(self, db: Session):
        self.db = db
        self.api_key = os.getenv("ANTHROPIC_API_KEY")

    def extract_text_from_file(self, file_path: str, content_type: str) -> str:
        if "presentation" in content_type or file_path.endswith((".pptx", ".ppt")):
            return self._extract_pptx(file_path)
        elif "pdf" in content_type or file_path.endswith(".pdf"):
            return self._extract_pdf(file_path)
        else:
            raise ValueError(f"Unsupported file type for analysis: {content_type}")

    def _extract_pptx(self, file_path: str) -> str:
        if not HAS_PPTX:
            raise ImportError("python-pptx required. Run: pip install python-pptx")
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            texts.append(cell.text)
        return "\n".join(texts)

    def _extract_pdf(self, file_path: str) -> str:
        if not HAS_PDF:
            raise ImportError("PyPDF2 required. Run: pip install PyPDF2")
        reader = PdfReader(file_path)
        texts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                texts.append(text)
        return "\n".join(texts)

    async def analyze_themes(self, text: str) -> List[str]:
        if not self.api_key:
            raise ValueError(
                "ANTHROPIC_API_KEY not configured. "
                "Add it to backend/.env to enable AI presentation analysis."
            )

        truncated = text[:8000]
        all_tags = self._get_all_relevance_tags()
        tag_context = ", ".join(sorted(all_tags)[:100]) if all_tags else "No existing tags yet"

        prompt = (
            "Analyze this presentation text and extract the key business themes, "
            "topics, sectors, and capabilities it covers. Return them as a JSON array of short tag strings.\n\n"
            f"Consider these existing tags in the system for consistency:\n{tag_context}\n\n"
            "Map the presentation content to relevant tags. Include:\n"
            "- Sector tags (e.g., 'Public Sector', 'Insurance', 'Banking')\n"
            "- Capability tags (e.g., 'Audit', 'Risk Management', 'Compliance')\n"
            "- Topic tags (e.g., 'Solvency II', 'IFRS 9', 'ESG')\n"
            "- Role relevance tags (e.g., 'Leadership', 'Decision-Maker')\n"
            "- Geography tags if mentioned (e.g., 'Nordics', 'EU')\n\n"
            f"Presentation text:\n{truncated}\n\n"
            "Return ONLY a JSON array of tag strings. "
            'Example: ["Insurance", "Risk Management", "Solvency II", "Nordics"]'
        )

        async with httpx.AsyncClient() as client:
            resp = await client.post(
                "https://api.anthropic.com/v1/messages",
                headers={
                    "x-api-key": self.api_key,
                    "content-type": "application/json",
                    "anthropic-version": "2023-06-01",
                },
                json={
                    "model": "claude-sonnet-4-20250514",
                    "max_tokens": 1024,
                    "messages": [{"role": "user", "content": prompt}],
                },
                timeout=30.0,
            )
            resp.raise_for_status()
            data = resp.json()
            text_response = data["content"][0]["text"]

            cleaned = text_response.strip()
            if cleaned.startswith("```"):
                lines = cleaned.split("\n")
                cleaned = "\n".join(lines[1:-1] if lines[-1].strip() == "```" else lines[1:])
            return json.loads(cleaned)

    def _get_all_relevance_tags(self) -> set:
        raw_tags = self.db.query(Contact.relevance_tags).filter(
            Contact.relevance_tags.isnot(None)
        ).distinct().all()
        tags = set()
        for (raw,) in raw_tags:
            try:
                parsed = json.loads(raw)
                if isinstance(parsed, list):
                    tags.update(t.strip() for t in parsed if isinstance(t, str) and t.strip())
            except (json.JSONDecodeError, TypeError):
                tags.update(t.strip() for t in raw.split(",") if t.strip())
        return tags

    def find_matching_contacts(self, themes: List[str], limit: int = 100) -> List[dict]:
        themes_lower = [t.lower() for t in themes]

        contacts = (
            self.db.query(Contact)
            .filter(
                Contact.status == ContactStatusEnum.ACTIVE,
                Contact.relevance_tags.isnot(None),
                Contact.email.isnot(None),
                Contact.email != "",
            )
            .all()
        )

        suppressed_ids = set(
            r[0] for r in self.db.query(SuppressionEntry.contact_id)
            .filter(SuppressionEntry.is_active == True).all()
        )

        scored = []
        for c in contacts:
            if c.id in suppressed_ids:
                continue
            try:
                c_tags = json.loads(c.relevance_tags)
            except (json.JSONDecodeError, TypeError):
                c_tags = [t.strip() for t in c.relevance_tags.split(",") if t.strip()]

            c_tags_lower = [t.lower() for t in c_tags]
            matched = []
            for theme in themes_lower:
                for tag in c_tags_lower:
                    if theme == tag or theme in tag or tag in theme:
                        matched.append(tag)
                        break

            if matched:
                score = len(matched)
                if c.is_decision_maker:
                    score += 0.5
                scored.append({
                    "contact_id": c.id,
                    "full_name": c.full_name,
                    "email": c.email,
                    "company_name": c.company_name,
                    "job_title": c.job_title,
                    "relevance_tags": c.relevance_tags,
                    "match_score": score,
                    "matched_tags": list(set(matched)),
                    "is_decision_maker": bool(c.is_decision_maker),
                })

        scored.sort(key=lambda x: x["match_score"], reverse=True)
        return scored[:limit]
